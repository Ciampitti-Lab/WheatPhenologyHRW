"""Build the defense-grade presentation for the WheatPhenologyHRW paper.

Faithful to the current manuscript (paper-overleaf): 7-model Phase-E with
FT-Transformer adopted at anthesis/maturity, honest 5/8 best-vs-best with
the controlled with/without-WES ablation as primary evidence, DOS target,
maturity correction + caveat. Embeds the regenerated figures F1-F6.

Run:  python scripts/build_presentation.py
Out:  outputs/WheatPhenologyHRW_Presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
FIG  = ROOT / 'docs' / 'figures'
OUT  = ROOT / 'outputs' / 'WheatPhenologyHRW_Presentation.pptx'
OUT.parent.mkdir(parents=True, exist_ok=True)

# Purdue palette
GOLD   = RGBColor(0xCE, 0xB8, 0x88)
ACCENT = RGBColor(0x8E, 0x6F, 0x3E)
DARK   = RGBColor(0x1B, 0x1B, 0x1B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF7, 0xF5, 0xEF)
GREY   = RGBColor(0x5A, 0x5A, 0x5A)
GREEN  = RGBColor(0x1B, 0x6E, 0x63)
RED    = RGBColor(0xA0, 0x39, 0x39)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height


def _blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = LIGHT
    return s

def _box(s, x, y, w, h, fill=None, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(1)
    return sh

def _txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    if isinstance(runs, str): runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        if isinstance(para, tuple): para = [para]
        for text, st in para:
            r = p.add_run(); r.text = text
            r.font.size = Pt(st.get('sz', 18)); r.font.bold = st.get('b', False)
            r.font.italic = st.get('i', False)
            r.font.color.rgb = st.get('c', DARK)
            r.font.name = 'Calibri'
    return tb

def _header(s, title, kicker=None):
    _box(s, 0, 0, SW, Inches(1.15), fill=DARK)
    _box(s, 0, Inches(1.15), SW, Pt(4), fill=GOLD)
    if kicker:
        _txt(s, Inches(0.55), Inches(0.13), Inches(12), Inches(0.3),
             [[(kicker.upper(), {'sz': 12, 'b': True, 'c': GOLD})]])
    _txt(s, Inches(0.55), Inches(0.36 if kicker else 0.28),
         Inches(12.2), Inches(0.75),
         [[(title, {'sz': 26, 'b': True, 'c': WHITE})]], anchor=MSO_ANCHOR.MIDDLE)

def _pagenum():
    for i, s in enumerate(prs.slides, 1):
        if i == 1: continue
        _txt(s, SW - Inches(0.9), SH - Inches(0.42), Inches(0.7), Inches(0.3),
             [[(str(i), {'sz': 11, 'c': GREY})]], align=PP_ALIGN.RIGHT)

def bullets(title, items, kicker=None, foot=None):
    s = _blank(); _header(s, title, kicker)
    y = Inches(1.5)
    for it in items:
        lvl = it[0] if isinstance(it, tuple) else 0
        text = it[1] if isinstance(it, tuple) else it
        bx = Inches(0.7 + 0.5 * lvl)
        mk = '▪' if lvl == 0 else '–'
        _txt(s, bx, y, Inches(0.35), Inches(0.4),
             [[(mk, {'sz': 18 if lvl else 20, 'b': True, 'c': ACCENT})]])
        _txt(s, bx + Inches(0.42), y, Inches(11.6 - 0.5 * lvl), Inches(0.9),
             [[(text, {'sz': 19 - 2 * lvl, 'b': (lvl == 0 and text.endswith(':')),
                       'c': DARK})]], line_spacing=1.0)
        y += Inches(0.62 if lvl == 0 else 0.5)
    if foot:
        _box(s, 0, SH - Inches(0.7), SW, Inches(0.7), fill=GOLD)
        _txt(s, Inches(0.55), SH - Inches(0.62), Inches(12.2), Inches(0.55),
             [[(foot, {'sz': 15, 'b': True, 'c': DARK})]], anchor=MSO_ANCHOR.MIDDLE)
    return s

def image_slide(title, img, kicker=None, caption=None, side=None):
    s = _blank(); _header(s, title, kicker)
    p = FIG / img
    iw = Inches(8.6 if side else 12.0)
    try:
        from PIL import Image
        w0, h0 = Image.open(p).size; ar = h0 / w0
    except Exception:
        ar = 0.52
    max_h = Inches(5.0)
    if iw * ar > max_h: iw = Emu(int(max_h / ar))
    x = Inches(0.55) if side else (SW - iw) // 2
    y = Inches(1.45)
    s.shapes.add_picture(str(p), x, y, width=iw)
    if side:
        _txt(s, Inches(9.4), Inches(1.7), Inches(3.5), Inches(5.0),
             [[(b, {'sz': 16, 'c': DARK})] for b in side], line_spacing=1.05,
             space_after=10)
    if caption:
        _txt(s, Inches(0.55), SH - Inches(0.62), Inches(12.2), Inches(0.5),
             [[(caption, {'sz': 13, 'i': True, 'c': GREY})]], align=PP_ALIGN.CENTER)
    return s

def table_slide(title, headers, rows, kicker=None, foot=None, hi=None):
    s = _blank(); _header(s, title, kicker)
    nr, nc = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nr, nc, Inches(0.55), Inches(1.5),
                            Inches(12.2), Inches(0.42 * nr)).table
    for j, htxt in enumerate(headers):
        c = gt.cell(0, j); c.text = htxt
        c.fill.solid(); c.fill.fore_color.rgb = DARK
        pr = c.text_frame.paragraphs[0]; pr.alignment = PP_ALIGN.CENTER
        pr.runs[0].font.size = Pt(15); pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            pr = c.text_frame.paragraphs[0]
            pr.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            pr.runs[0].font.size = Pt(14)
            pr.runs[0].font.color.rgb = DARK
            if hi and (i - 1) in hi:
                c.fill.fore_color.rgb = GOLD
                pr.runs[0].font.bold = True
    if foot:
        _txt(s, Inches(0.55), SH - Inches(0.85), Inches(12.2), Inches(0.7),
             [[(foot, {'sz': 14, 'i': True, 'c': GREY})]])
    return s

def section(title, sub=None):
    s = _blank()
    _box(s, 0, 0, SW, SH, fill=DARK)
    _box(s, 0, Inches(3.05), SW, Pt(4), fill=GOLD)
    _txt(s, Inches(0.9), Inches(2.95), Inches(11.5), Inches(1.1),
         [[(title, {'sz': 40, 'b': True, 'c': WHITE})]], anchor=MSO_ANCHOR.BOTTOM)
    if sub:
        _txt(s, Inches(0.95), Inches(4.2), Inches(11.5), Inches(0.6),
             [[(sub, {'sz': 18, 'c': GOLD})]])
    return s

# ───────────────────────── TITLE ─────────────────────────
s = _blank()
_box(s, 0, 0, SW, SH, fill=DARK)
_box(s, 0, Inches(4.5), SW, Pt(5), fill=GOLD)
_txt(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.6),
     [[("A multi-stage hybrid framework for", {'sz': 32, 'b': True, 'c': WHITE})],
      [("field-scale winter-wheat phenology", {'sz': 32, 'b': True, 'c': WHITE})],
      [("across the U.S. Hard Red Winter belt", {'sz': 32, 'b': True, 'c': GOLD})]],
     line_spacing=1.05)
_txt(s, Inches(0.85), Inches(4.75), Inches(11.7), Inches(1.4),
     [[("Vlasios Mangidis · German Mandrini · Ignacio A. Ciampitti", {'sz': 18, 'c': WHITE})],
      [("Department of Agronomy, Purdue University", {'sz': 15, 'c': GREY})],
      [("Eight stages · 5,293 fields · 8,465 field-years · HLS + Daymet + MODIS-LST + Wang–Engel–Streck",
        {'sz': 14, 'i': True, 'c': GOLD})]], space_after=8)

bullets("Outline", [
    "Motivation: why field-scale winter-wheat phenology",
    "The gap in the literature & our contribution",
    "Data: study area, ground labels, satellite & meteorology",
    "Framework: feature engineering + Wang–Engel–Streck predictor",
    "Modelling: 7 candidate models, 2 strategies, LOYO/LOSO",
    "Results: accuracy, strategy value, robustness, importance, transferability",
    "Honest delineation: where it works, where it does not",
    "Discussion, limitations, conclusion — backup detail",
], kicker="Roadmap")

# ───────────────────────── MOTIVATION ─────────────────────────
section("1 · Motivation & Gap")

bullets("Why field-scale phenology timing matters", [
    "HRW belt = largest U.S. wheat production by class",
    "Stage timing governs yield & quality formation:",
    (1, "N & disease management, irrigation, frost-risk, harvest planning"),
    (1, "insurance phenotyping, post-season yield/protein analytics"),
    "Operational source (USDA-NASS Crop Progress): state-level, weekly",
    (1, "masks within-state heterogeneity: sowing, cultivar, soil, elevation, microclimate"),
], kicker="Problem",
   foot="Need: close the gap between regional cadence and field-scale decisions")

bullets("The eight phenological stages", [
    "Emergence → Tillering → Jointing → Flag leaf → Boot → Heading → Anthesis → Maturity",
    "Early vegetative (emergence/tillering/jointing): bare-soil regime, dormancy-interrupted, noisy descriptive labels",
    "Reproductive (flag leaf/boot/heading/anthesis): peak canopy, tightly thermally gated, unambiguous labels",
    "Maturity: small cleaned cohort, treated as preliminary",
], kicker="Target")

bullets("What no prior framework does — all four at once", [
    "(i) predicts the full 8-stage winter-wheat cycle",
    "(ii) operates at field scale across a continental climate gradient",
    "(iii) explicitly fuses physiology-based thermal-time with ML corrections",
    "(iv) validates spatial transferability via leave-one-state-out CV",
    "Prior work is narrower on one or more axes (single stage / region / pixel-level)",
], kicker="Literature gap",
   foot="Novelty is the combination — not a single new estimator")

bullets("Contributions", [
    "First full 8-stage, field-scale framework across the HRW belt (~10° latitude, TX→CO/NE)",
    "Physiology + satellite fusion: Wang–Engel–Streck predictor anchored to per-field OBSERVED sowing",
    "Observed-sowing anchor: removes a calendar confound; effect quantified, not assumed",
    "Explicit spatial-transferability test (LOSO) — reported honestly, including where it breaks",
    "Reproducible: public code, GEE asset, end-to-end pipeline",
], kicker="What we deliver")

# ───────────────────────── DATA ─────────────────────────
section("2 · Data & Study Area")

image_slide("Study area: the U.S. Hard Red Winter belt", "F1_study_area.png",
            kicker="Dataset", side=[
            "5,293 winter-wheat fields",
            "8,465 field-years",
            "4 growing seasons 2013/14–2016/17",
            "Kansas-dominated (~72% of fields)",
            "Strong gradients: temperature, precipitation, elevation",
            "TX → OK → KS → NE → CO"],
            caption="Per-field ground phenology across a continental climate gradient")

bullets("Ground phenology labels & the DOS target", [
    "Per (field, harvest-year): earliest day a stage-vocabulary label is recorded",
    "Target = Day of Season (DOS): days since 1 July of harvest-year−1",
    (1, "spans the full season (prev. summer → current summer harvest)"),
    "Maturity quality filter: discard DOS < 280 (previous-season harvest residue)",
    (1, "removed 295 of 788 raw maturity obs (37%) → cohort 357 field-years"),
    "Descriptive (non-Zadoks) labels → intrinsic scatter, bounded accuracy",
], kicker="Target variable")

table_slide("Satellite & meteorological inputs",
            ["Source", "Variable", "Role"],
            [["Harmonized Landsat–Sentinel", "Surface reflectance, 30 m, 2–3 d", "VI time series → phenometrics"],
             ["Daymet", "Daily T, precip, radiation", "Meteorology + thermal-time"],
             ["MODIS", "Land-surface temperature", "Canopy thermal state"],
             ["Wang–Engel–Streck", "Thermal-time stage predictor", "Physiology prior (Hybrid)"]],
            kicker="Inputs",
            foot="All features from the full season: 1 Jul (y−1) → 31 Jul (y); same vector predicts every stage")

# ───────────────────────── FRAMEWORK ─────────────────────────
section("3 · Framework & Modelling")

image_slide("End-to-end framework", "F2_framework.png", kicker="Pipeline",
            caption="Satellite phenometrics + Daymet + MODIS-LST + WES predictor → per-stage ML → DOS")

bullets("Feature engineering", [
    "Full-season VI time series (HLS): smoothed (interp + Savitzky–Golay)",
    "Phenometric descriptors of curve shape & timing",
    "Daymet meteorology + thermal-time accumulators",
    "MODIS LST canopy thermal features; state encoders",
    "Feature groups: WES · HLS phenometrics · MODIS LST · Thermal-time · Daymet meteo · State",
], kicker="Inputs → features")

bullets("The Wang–Engel–Streck thermal-time predictor", [
    "Process-based engine: cardinal-temperature + generalised vernalisation + photoperiod",
    "Models cold/daylength gating specific to winter wheat (vs photo-thermal accumulator alone)",
    "Anchored to per-field OBSERVED sowing where available:",
    (1, "16.3% observed anchor · 83.7% state-median fallback"),
    "Predicts canonical stage sequence → used as ML feature in the Hybrid strategy",
], kicker="Physiology core")

bullets("Two strategies", [
    "ML-only (B): satellite + meteo features, no WES",
    "Physiology-informed / Hybrid (C): WES features + ML",
    "Same feature vector predicts all 8 stages; model learns per-stage relevance",
    "Strategy comparison = a feature-set question (value of the physiology prior)",
], kicker="Design", foot="Hybrid lets ML refine, weight or ignore the WES predictor")

bullets("Seven candidate models per stage × strategy", [
    "Linear: ElasticNet, Ridge",
    "Trees: Random Forest, XGBoost, LightGBM",
    "Deep tabular: TabNet, FT-Transformer",
    "Deep models evaluated fairly (standardised I/O, inner-year early stopping, 5-seed avg)",
    "Selected per stage strictly by held-out R²",
], kicker="Modelling",
   foot="We did NOT exclude deep models a priori — we tested them under the identical protocol")

bullets("Cross-validation", [
    "Leave-One-Year-Out (LOYO): temporal generalisation, primary accuracy metric",
    "Leave-One-State-Out (LOSO): spatial transferability under spatial autocorrelation",
    (1, "withholds each state's data AND its encoder entirely — a strict test"),
    "Bootstrap 95% CIs (2,000 resamples) on every reported R²",
], kicker="Evaluation")

bullets("Phase-E experimental grid", [
    "8 stages × 2 strategies × 7 models = full grid (Supplement S1)",
    "Linear models keep an inner SelectKBest tuning loop",
    "Best (strategy, model) per stage → carried to the headline table",
    "Ties at 2 decimals resolved toward the interpretable model",
], kicker="How 'best' is chosen")

# ───────────────────────── RESULTS ─────────────────────────
section("4 · Results")

image_slide("Per-stage prediction accuracy (LOYO)", "F3_per_stage_scatter.png",
            kicker="Accuracy",
            caption="Reproductive R² 0.69–0.82 (RMSE 4.6–5.8 d) · early vegetative 0.33–0.36 · maturity 0.44 (preliminary)")

table_slide("Best model per stage",
            ["Stage", "Strategy", "Model", "R²", "RMSE (d)", "n"],
            [["Emergence", "Physiology", "LightGBM", "0.36", "29.0", "4,145"],
             ["Tillering", "ML-only", "ElasticNet", "0.34", "16.5", "4,715"],
             ["Jointing", "Physiology", "LightGBM", "0.33", "16.2", "5,124"],
             ["Flag leaf", "Physiology", "XGBoost", "0.71", "5.8", "2,657"],
             ["Boot", "Physiology", "LightGBM", "0.69", "5.4", "2,320"],
             ["Heading", "Physiology", "ElasticNet", "0.73", "5.5", "2,714"],
             ["Anthesis", "Physiology", "FT-Transformer", "0.82", "4.6", "1,464"],
             ["Maturity", "Physiology", "FT-Transformer", "0.44", "5.7", "357"]],
            kicker="Headline numbers", hi=[6, 7],
            foot="FT-Transformer adopted only where it genuinely wins (anthesis, maturity)")

bullets("Reproductive vs vegetative skill gap", [
    "Reproductive cluster: peak canopy biomass, sharp thermal gating, unambiguous labels",
    "Early vegetative: bare-soil signal, dormancy-interrupted window, coarse cadence, noisy labels",
    "Gap ≈ 0.40 in R² (and 10–23 d RMSE) — set substantially by the GROUND DATA",
    "Empirical label scatter: tillering 8.8 d, jointing 8.1 d, emergence 6.8 d (S1)",
], kicker="The salient pattern",
   foot="The early-stage ceiling is a data noise floor, not a model deficit")

bullets("The maturity correction — reported honestly", [
    "Earlier build retained previous-season 'Harvest Ready' residue → bimodal, inflated R²=0.74",
    "Removing DOS<280 collapses the inflation; corrected cohort n=357",
    "Best model (FT-Transformer) reaches R²=0.44 — still weakest & most uncertain",
    "Caveat: at maturity ΔWES ≈ 0 — the FT gain is meteo capacity, NOT the physiology prior",
], kicker="Evaluation discipline",
   foot="Surfaced in the main text, not buried — the clearest illustration of the discipline")

image_slide("Feature-strategy comparison", "F4_strategy_comparison.png",
            kicker="Value of physiology",
            caption="Best-model-vs-best-model: physiology-informed higher in 5/8 (2 ties, tillering ML-only)")

bullets("Why 5/8 — not 7/8", [
    "5/8 = best-of-7-with-WES vs best-of-7-without-WES, at 2 decimals",
    "emergence & maturity ΔR² = +0.001 → round to 0.00 AND within bootstrap CI = ties",
    "tillering −0.04 (gradual, no sharp thermal trigger; sowing-anchor noise dominates)",
    "We do NOT count within-noise deltas as wins — conservative, defensible",
    "A strong FT-Transformer is available to BOTH strategies → best-vs-best is a lower bound",
], kicker="Scientific conservatism")

table_slide("Controlled with/without-WES ablation — the robust evidence",
            ["Stage", "Control gain", "σ=7 d", "σ=14 d", "σ=21 d", "σ=28 d"],
            [["Flag leaf", "0.274", "82%", "81%", "78%", "57%"],
             ["Boot", "0.120", "89%", "87%", "83%", "76%"],
             ["Heading", "0.336", "94%", "94%", "88%", "88%"],
             ["Anthesis (FT)", "0.025", "105%", "108%", "90%", "119%"]],
            kicker="Primary physiology evidence",
            foot="Same model, WES on/off (model-agnostic). Flag-leaf/boot/heading: large gains, "
                 "robust to σ=28 d (> real sowing spread). Anthesis: FT re-learns the signal → small "
                 "gain (0.025) but still robust → physiology contribution is real, not leakage")

image_slide("Feature importance across stages", "F5_feature_importance.png",
            kicker="Mechanism",
            caption="WES dominates flag leaf (~73%); reproductive shifts toward Daymet meteo; deep-stage = permutation importance")

bullets("Deep-model finding (Phase-E S1, 8×2×7)", [
    "TabNet: unstable, never competitive at these cohort sizes (350–5,100 field-years)",
    "FT-Transformer: competitive at reproductive stages; best at anthesis (0.82) & maturity (0.44)",
    "Trees/linear remain best at the other 6 stages — kept for per-feature interpretability",
    "Honest consequence: with strong FT in ML-only too, the best-vs-best physiology gap narrows",
], kicker="Deep tabular models",
   foot="Best-vs-best narrows because FT re-learns WES — see next slide")

bullets("Why the physiology prior (WES) matters", [
    "The small best-vs-best ΔR² is NOT evidence WES is redundant",
    "Controlled ablation (model held fixed): WES adds +0.34 R² (heading), +0.27 (flag leaf), +0.12 (boot)",
    "A deep model in ML-only narrows the gap only by re-learning — opaquely, imperfectly — the thermal-time relation WES supplies explicitly",
    "→ WES lets a simple, interpretable tree/linear model reach deep-model accuracy",
    "→ on a mechanistic, transferable basis (vernalisation, photoperiod, cardinal temps) a black-box lacks",
    "Robustness: only trees + linear stable at every stage; TabNet collapses → headline not on a fragile estimator",
], kicker="The core argument",
   foot="Best-vs-best understates WES; the model-fixed ablation is the true measure")

image_slide("Spatial transferability (LOSO)", "F6_loso_transferability.png",
            kicker="Transferability",
            caption="Moderate across southern/central states at reproductive stages; Colorado consistently hardest")

bullets("Reading the LOSO result", [
    "Strict test: each state's data + encoder withheld entirely",
    "Kansas-dominated model transfers with bounded skill to adjacent southern belt",
    "Colorado weakest: high-plains elevation, later cold-tolerant cultivars, dispersed sowing",
    "Scope delineation, not a defect — field-scale resolves heterogeneity a regional mean masks",
    "Negative control (S8): controlled LOSO with/without-WES → WES does NOT aid cross-region transfer (neutral at well-sampled states); it is an in-distribution thermal-time prior, not a spatial-generalisation device — reported honestly",
], kicker="Interpretation")

# ───────────────────────── DISCUSSION ─────────────────────────
section("5 · Discussion & Conclusion")

bullets("Selective value of the physiology prior", [
    "WES dominates flag leaf — a sharp thermally-gated transition (cardinal-temperature response)",
    "Contributes modestly where photoperiod/vernalisation/soil-moisture mean no single engine suffices",
    "Tillering: only ML-only win — gradual, no sharp trigger; anchor uncertainty > thermal signal",
    "Extends the PhenoCrop satellite+process design to cold/daylength-gated winter wheat",
], kicker="Why it works where it works")

bullets("Positioning vs prior work", [
    "Comparable accuracy on the most comparable target (field-scale anthesis ≈ 0.82)",
    "MODIS/Landsat fusion & HLS planting-date work: similar errors, but single-crop/single-stage",
    "Double-logistic NDVI baseline ≈ unskilled at heading/anthesis (R²≈0.02) — signal not in curve shape alone",
    "No multi-state, field-scale benchmark exists for tillering/jointing in winter wheat",
], kicker="Novelty in context")

bullets("Limitations", [
    "Only 4 seasons — climatic extremes (drought, late-freeze) untested",
    "Descriptive non-Zadoks labels → irreducible noise; tillering partly a target-definition floor",
    "Maturity: smaller cleaned cohort, treated most cautiously",
    "State-median sowing fallback may misanchor atypical fields (bounded by σ-perturbation test)",
    "Retrospective (full-season) — not an in-season forecaster (productive follow-up)",
], kicker="Honest caveats")

bullets("Conclusion", [
    "First multi-stage, field-scale, physiology-informed winter-wheat phenology framework for the HRW belt",
    "Reproductive transitions R² 0.69–0.82 (RMSE 4.6–5.8 d); early stages honestly bounded by label noise",
    "Physiology contribution robust under a controlled, model-agnostic with/without-WES ablation",
    "Transferability quantified — strong south/central, bounded at the high-elevation margin",
    "Supports retrospective label generation to train the short-horizon forecasters hand data cannot",
], kicker="Take-home",
   foot="Reproducible field-scale basis for HRW-wheat phenology — design transfers with recalibration")

bullets("Future work", [
    "In-season variant: retrain on truncated feature windows for early warning",
    "Zadoks-consistent re-labelling study to pin the vegetative noise floor",
    "Deterministic observed-anchor→median refit (confirmatory systematic-confound test)",
    "Recalibration to Soft Red Winter / European winter-wheat systems",
], kicker="Next")

# ───────────────────────── BACKUP ─────────────────────────
section("Backup", "Supplementary detail")

table_slide("S1 · Phase-E grid — selected reproductive rows (Hybrid R²)",
            ["Stage", "ElasticNet", "XGBoost", "LightGBM", "TabNet", "FT-Transformer"],
            [["Flag leaf", "0.66", "0.71*", "0.70", "0.65", "0.71"],
             ["Boot", "0.64", "0.68", "0.69*", "0.57", "0.68"],
             ["Heading", "0.73*", "0.69", "0.71", "−1.72", "0.73"],
             ["Anthesis", "0.80", "0.71", "0.70", "0.68", "0.82*"],
             ["Maturity", "0.34", "−0.06", "0.34", "−2.79", "0.44*"]],
            kicker="Full grid in Supplement S1",
            foot="* = carried to main table. FT ties trees/linear at flag leaf/heading → interpretable model kept")

bullets("S3 · Maturity label correction", [
    "Raw 'Maturity'/'Harvest Ready' includes previous-season residue (physiologically impossible)",
    "Filter DOS < 280 → remove 295 of 788 raw obs (37%)",
    "Bimodal target collapses; honest R² = 0.44 on 357 field-years (was inflated 0.74)",
], kicker="Backup")

bullets("S4 · Sowing-anchor robustness method", [
    "Model fixed per stage; add Gaussian noise σ∈{7,14,21,28} d to fallback anchors only",
    "Re-run WES, recompute controlled gain; 3 replicates each",
    "Calibrated to reality: observed sowing departs state median by 14.7 d MAD (σ=21.5 d)",
    "σ=28 d is a conservative worst case — exceeds empirical state-median error",
], kicker="Backup")

bullets("S6/S7 · Baselines & tillering floor", [
    "Double-logistic NDVI baseline: R²≈0.02 at heading/anthesis (signal not in curve shape)",
    "Zadoks-consolidated tillering target does not recover skill → ceiling is label/signal, not granularity",
    "Label-intrinsic onset scatter: tillering 8.8 d, jointing 8.1 d, emergence 6.8 d",
], kicker="Backup")

bullets("Reproducibility", [
    "Public code repository (end-to-end pipeline: features → Phase-E → figures)",
    "Google Earth Engine field-buffer asset",
    "Bootstrap CIs and full Phase-E grid (S1) reported",
    "Manuscript: paper-overleaf (target: Int. J. Applied Earth Obs. Geoinformation)",
], kicker="Open science",
   foot="Questions?")

_pagenum()
prs.save(str(OUT))
print(f"Saved {OUT} — {len(prs.slides)} slides")
